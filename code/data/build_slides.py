#!/usr/bin/env python3
"""Build the AI Planner term-project deck (English, ~15 slides, 16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INDIGO = RGBColor(0x26, 0x3A, 0x6B)
INDIGO2 = RGBColor(0x3B, 0x57, 0x9E)
GOLD = RGBColor(0xE0, 0xA1, 0x32)
DARK = RGBColor(0x1E, 0x1E, 0x24)
GRAY = RGBColor(0x55, 0x5A, 0x66)
LIGHT = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FOOT = "Han Seo-hee · GIST AI Convergence — Building a Useful Private AI Butler · 2026"


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def fill(slide, color):
    s = slide.shapes.add_shape(1, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def band(slide, t, h, color):
    s = slide.shapes.add_shape(1, 0, t, SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def footer(slide, idx):
    tb = _box(slide, Inches(0.5), Inches(7.05), Inches(11), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = FOOT
    r.font.size = Pt(9); r.font.color.rgb = GRAY
    n = _box(slide, Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4))
    pn = n.text_frame.paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run(); rn.text = str(idx)
    rn.font.size = Pt(11); rn.font.color.rgb = INDIGO2; rn.font.bold = True


def image_slide(path):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(path, 0, 0, width=SW, height=SH)
    return s


def title_slide():
    s = prs.slides.add_slide(BLANK)
    fill(s, INDIGO)
    band(s, Inches(3.05), Inches(0.06), GOLD)
    tb = _box(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.7))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "From a Flashy Butler to a Useful One"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    sub = _box(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.0))
    p2 = sub.text_frame.paragraphs[0]
    r2 = p2.add_run(); r2.text = "A Cost-Bounded Weekly Research Coach for Grad Students"
    r2.font.size = Pt(22); r2.font.color.rgb = GOLD; r2.font.italic = True
    meta = _box(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.2))
    for i, line in enumerate([
        "Han Seo-hee  ·  GIST, AI Convergence",
        "Term Project — Personal AI Butler  ·  June 2026",
    ]):
        para = meta.text_frame.paragraphs[0] if i == 0 else meta.text_frame.add_paragraph()
        rr = para.add_run(); rr.text = line
        rr.font.size = Pt(16); rr.font.color.rgb = WHITE


def content_slide(idx, title, kicker, bullets):
    s = prs.slides.add_slide(BLANK)
    fill(s, WHITE)
    band(s, 0, Inches(1.35), INDIGO)
    band(s, Inches(1.35), Inches(0.05), GOLD)
    # kicker
    kb = _box(s, Inches(0.6), Inches(0.22), Inches(11), Inches(0.4))
    kp = kb.text_frame.paragraphs[0]
    kr = kp.add_run(); kr.text = kicker.upper()
    kr.font.size = Pt(12); kr.font.bold = True; kr.font.color.rgb = GOLD
    # title
    tb = _box(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.8))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run(); tr.text = title
    tr.font.size = Pt(30); tr.font.bold = True; tr.font.color.rgb = WHITE
    # body
    body = _box(s, Inches(0.75), Inches(1.75), Inches(11.9), Inches(5.0))
    tf = body.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        lvl = 0; text = b
        if isinstance(b, tuple):
            lvl, text = b
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(8)
        bold = text.startswith("**")
        text = text.replace("**", "")
        run = p.add_run(); run.text = ("•  " if lvl == 0 else "–  ") + text
        run.font.size = Pt(20 - 3 * lvl)
        run.font.color.rgb = DARK if lvl == 0 else GRAY
        if bold:
            run.font.bold = True; run.font.color.rgb = INDIGO2
    footer(s, idx)


def closing_slide(idx, title, lines):
    s = prs.slides.add_slide(BLANK)
    fill(s, INDIGO)
    band(s, Inches(2.4), Inches(0.06), GOLD)
    tb = _box(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = WHITE
    body = _box(s, Inches(0.95), Inches(2.7), Inches(11.4), Inches(4.0))
    tf = body.text_frame; tf.word_wrap = True
    for i, b in enumerate(lines):
        lvl = 0; text = b
        if isinstance(b, tuple):
            lvl, text = b
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(10)
        bold = text.startswith("**"); text = text.replace("**", "")
        run = p.add_run(); run.text = ("•  " if lvl == 0 else "–  ") + text
        run.font.size = Pt(20 - 2 * lvl)
        run.font.color.rgb = (GOLD if bold else WHITE)
        run.font.bold = bold
    footer(s, idx)


title_slide()

content_slide(2, "The Problem", "Motivation", [
    "**Grad students rarely know if they are on track.**",
    "Advisors see only part of the picture — you can't tell them everything.",
    "Every lab and every student differs → generic planners don't fit.",
    "**Missing: a private, honest signal of \"am I doing well this week?\"**",
])

content_slide(3, "Where I Started", "The Journey", [
    "An always-on Telegram butler: OpenClaw → Claude API → my machine.",
    "It could send email, post to my blog, log my music, write Obsidian notes, send daily briefings.",
    "**The first time I sent an email through Telegram, it felt magical.**",
])

content_slide(4, "But Magic ≠ Value", "The Turn", [
    "OpenClaw is **API-metered, not a subscription** — every single action costs money.",
    "Sending the email myself was cheaper and faster than paying the agent to do it.",
    "**Novelty fades; only usefulness justifies the spend.**",
    "If Claude or ChatGPT billed me per message, I'd use them far less, too.",
])

content_slide(5, "Subtraction as Design", "Design Principle", [
    "I removed everything that wasn't worth paying for:",
    (1, "Email, daily notifications, Obsidian logging (Linux already does it), music curation."),
    "A flat subscription would have hidden this — metered cost forced honesty.",
    "**What survived had to justify its cost, per run.**",
])

content_slide(6, "The Distilled Product", "What It Became", [
    "**A weekly research coach for grad students.**",
    "You log your week cheaply, yourself — essentially free.",
    "The AI synthesizes **once a week** (a few cents) into feedback grounded in YOUR week vs YOUR goals.",
    "Record Mon–Fri → get reviewed on the weekend.",
])

image_slide("/home/jyaenugu/diagram_architecture.png")

content_slide(8, "The Measuring Stick (the moat)", "Grounding", [
    "**RESEARCH_GOALS** — 2-year roadmap, degree-credit tracker, this-season focus, PhD-abroad target.",
    "The review measures my actual week **against my own goals** — not generic \"study more\" advice.",
    "**A chatbot without my goals and my activity simply cannot produce this.**",
])

content_slide(9, "Demo — The Weekly Review", "Practicality", [
    "One run produces six grounded sections:",
    (1, "① How the week went (day-by-day narrative)"),
    (1, "② Progress vs goals   ③ Gaps / drift   ④ Money spent"),
    (1, "⑤ Next-week schedule (from my calendar)   ⑥ 1–2 concrete focuses"),
    "**Grounded in real logs:** DAE paper, the lab-meeting talk, DDPM next steps — even sleep and spending.",
])

content_slide(10, "Cost Economy", "$ / month", [
    "**Two-tier cost:** logging = $0 (local) + synthesis ≈ a few cents / week (~tens of cents / month).",
    "**The $2 lesson:** an early full-context Sonnet run bled ~$2 in minutes.",
    "It confirmed the thesis: heavy always-on agents really do cost — **lean, on-demand wins.**",
])

content_slide(11, "Technical Rigour", "Systems", [
    "**Hard wall:** Tier-1 rate limit = 30,000 input tokens / minute (Sonnet).",
    "The agent loop resends full context + ~40 tool schemas every step → it blew past the limit.",
    "**Fix:** lightContext (drop persona bootstrap) + toolsAllow (6 tools, drop ~34 schemas) + model choice.",
    "Inference-cost framing: **I = M × HBM × R**.",
])

image_slide("/home/jyaenugu/diagram_dataflow.png")

content_slide(13, "Differentiation vs Big Tech", "≥ 3 of 5", [
    "**Grounded** in my real logged activity + my goals.",
    "**On my machine** — I own the data.",
    "**Cost-transparent** and cost-bounded.",
    "**Subtractive & personal** — shaped to one user, one lab.",
])

content_slide(14, "Smartening — RAG on My Own Vault", "Week-11 Method", [
    "The review is **retrieval-augmented generation over my own Vault.**",
    "It retrieves my notes, papers, and goals (read_note, recent_notes, recent_logs) to ground the feedback.",
    "**Week-11 method working:** retrieval-grounded generation on personal data.",
])

closing_slide(15, "Future Work & Takeaway", [
    "**v2 — Interactive check-in:** a 5-question Telegram dialogue before the review (catches what logs miss).",
    "**Then — other students:** prove it on myself first; they must log too.",
    "",
    "**\"A project no one uses is worth nothing.\"**",
    "**Built for daily-use value — not novelty.**",
])

out = "/home/jyaenugu/AI_Planner_GIST.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
